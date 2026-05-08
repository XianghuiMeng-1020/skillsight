#!/usr/bin/env python3
"""Build the SkillSight HKU Career Centre demo deck.

Run:
  .venv/bin/python scripts/build_demo_deck.py

Output:
  demo/SkillSight_HKU_CareerCenter.pptx
"""
from __future__ import annotations

from pathlib import Path
from io import BytesIO

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree

# ---------- Configuration ----------
ROOT = Path(__file__).resolve().parents[1]
SHOTS = ROOT / "demo" / "screenshots"
ASSETS = ROOT / "demo" / "assets"
OUT_PATH = ROOT / "demo" / "SkillSight_HKU_CareerCenter.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Brand palette
INK = RGBColor(0x1C, 0x19, 0x17)
INK_SOFT = RGBColor(0x44, 0x40, 0x3D)
MUTED = RGBColor(0x8A, 0x82, 0x7C)
PAPER = RGBColor(0xFF, 0xF9, 0xF0)
PAPER_DEEP = RGBColor(0xFA, 0xF1, 0xE2)
CORAL = RGBColor(0xE1, 0x81, 0x82)
AMBER = RGBColor(0xF9, 0xCE, 0x9C)
SKY = RGBColor(0xC9, 0xDD, 0xE3)
SAGE = RGBColor(0xBB, 0xCF, 0xC3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT_DISPLAY = "Calibri"
FONT_BODY = "Calibri"

# ---------- Helpers ----------

def add_full_bg(slide, color: RGBColor) -> None:
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False


def add_text(
    slide,
    text: str,
    *,
    left,
    top,
    width,
    height,
    font_name: str = FONT_BODY,
    size: int = 18,
    bold: bool = False,
    color: RGBColor = INK,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
    line_spacing: float | None = None,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_bullets(slide, items: list[str], *, left, top, width, height, size: int = 16, color: RGBColor = INK_SOFT, line_spacing: float = 1.25, bullet_color: RGBColor = CORAL):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(4)
        # Bullet marker run
        marker = p.add_run()
        marker.text = "■  "
        marker.font.name = FONT_BODY
        marker.font.size = Pt(size)
        marker.font.color.rgb = bullet_color
        marker.font.bold = True
        # Body run
        body = p.add_run()
        body.text = item
        body.font.name = FONT_BODY
        body.font.size = Pt(size)
        body.font.color.rgb = color
    return tb


def add_chip(slide, text: str, *, left, top, width=None, height=Inches(0.42), fill: RGBColor = PAPER_DEEP, color: RGBColor = INK, size: int = 12, bold: bool = True):
    if width is None:
        width = Inches(max(1.5, 0.13 * len(text) + 0.3))
    chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    chip.adjustments[0] = 0.5
    chip.line.fill.background()
    chip.fill.solid()
    chip.fill.fore_color.rgb = fill
    tf = chip.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = FONT_BODY
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return chip


def add_rounded_image(slide, image_path: Path, *, left, top, width, height, border_color: RGBColor = WHITE, border_pt: float = 4, shadow: bool = True):
    """Place an image into a fixed box, scaling to fit while preserving aspect.

    Image is centered within (left, top, width, height) and clipped to a thin
    white card with subtle shadow.
    """
    if not image_path.exists():
        # Placeholder
        ph = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        ph.adjustments[0] = 0.04
        ph.fill.solid(); ph.fill.fore_color.rgb = PAPER_DEEP
        ph.line.color.rgb = MUTED
        ph.line.width = Pt(0.75)
        tf = ph.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = f"[ {image_path.name} ]"; r.font.color.rgb = MUTED; r.font.size = Pt(14)
        return ph

    img = Image.open(image_path)
    iw, ih = img.size
    box_ratio = width / height
    img_ratio = iw / ih
    if img_ratio >= box_ratio:
        # image wider — limit by width
        draw_w = width
        draw_h = Emu(int(width * ih / iw))
    else:
        draw_h = height
        draw_w = Emu(int(height * iw / ih))
    draw_left = left + Emu(int((width - draw_w) / 2))
    draw_top = top + Emu(int((height - draw_h) / 2))

    # Card backdrop
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.adjustments[0] = 0.025
    card.fill.solid(); card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = border_color
    card.line.width = Pt(border_pt * 0.25)

    pic = slide.shapes.add_picture(str(image_path), draw_left, draw_top, width=draw_w, height=draw_h)
    return pic


def add_step_bar(slide, current: int, total: int, *, top, left=Inches(0.7), width=Inches(11.93)):
    seg_w = width / total
    for i in range(total):
        seg_left = left + Emu(int(i * seg_w))
        is_active = i + 1 <= current
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, seg_left + Inches(0.04), top, seg_w - Inches(0.08), Inches(0.08))
        bar.adjustments[0] = 0.5
        bar.line.fill.background()
        bar.fill.solid()
        bar.fill.fore_color.rgb = CORAL if is_active else PAPER_DEEP
    add_text(
        slide,
        f"Step {current} of {total}",
        left=left, top=top + Inches(0.18), width=width, height=Inches(0.3),
        size=11, color=MUTED, align=PP_ALIGN.CENTER,
    )


def add_footer(slide, page: int):
    add_text(
        slide,
        "SkillSight  ·  HKU Skills-to-Jobs Transparency System  ·  GEM Lab, Faculty of Education",
        left=Inches(0.7), top=Inches(7.05), width=Inches(10.5), height=Inches(0.3),
        size=10, color=MUTED,
    )
    add_text(
        slide,
        f"{page:02d}",
        left=Inches(12.2), top=Inches(7.05), width=Inches(0.6), height=Inches(0.3),
        size=10, color=MUTED, align=PP_ALIGN.RIGHT,
    )


def add_brand_corner(slide):
    # subtle gradient corner block top-left
    block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.18), Inches(7.5))
    block.line.fill.background()
    block.fill.solid(); block.fill.fore_color.rgb = CORAL


# ---------- Slide layouts ----------

def slide_cover(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, PAPER)
    # decorative blobs
    blob1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(-1.5), Inches(8), Inches(8))
    blob1.line.fill.background(); blob1.fill.solid(); blob1.fill.fore_color.rgb = AMBER
    blob1.fill.fore_color.brightness = 0.0
    blob2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.5), Inches(2.5), Inches(8), Inches(8))
    blob2.line.fill.background(); blob2.fill.solid(); blob2.fill.fore_color.rgb = SKY
    # logo dot
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.95), Inches(0.95), Inches(0.45), Inches(0.45))
    dot.line.fill.background(); dot.fill.solid(); dot.fill.fore_color.rgb = CORAL
    add_text(slide, "SkillSight", left=Inches(1.5), top=Inches(0.92), width=Inches(4), height=Inches(0.55), size=22, bold=True, color=INK, font_name=FONT_DISPLAY)

    add_text(slide, "SkillSight", left=Inches(1.0), top=Inches(2.4), width=Inches(11), height=Inches(1.4), size=88, bold=True, color=INK, font_name=FONT_DISPLAY)
    add_text(slide, "Evidence-Based Career Readiness, for Every HKU Student", left=Inches(1.0), top=Inches(3.7), width=Inches(11), height=Inches(0.8), size=26, color=INK_SOFT, font_name=FONT_DISPLAY)
    # Divider line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(4.6), Inches(1.4), Inches(0.04))
    line.line.fill.background(); line.fill.solid(); line.fill.fore_color.rgb = CORAL
    add_text(slide, "A demonstration briefing for the HKU Centre of Development &", left=Inches(1.0), top=Inches(4.85), width=Inches(11), height=Inches(0.4), size=15, color=INK_SOFT)
    add_text(slide, "Resources for Students (CEDARS)  ·  April 2026", left=Inches(1.0), top=Inches(5.20), width=Inches(11), height=Inches(0.4), size=15, color=INK_SOFT)

    add_text(slide, "Developed by GEM Lab  ·  Faculty of Education, The University of Hong Kong",
             left=Inches(1.0), top=Inches(6.7), width=Inches(11), height=Inches(0.4), size=12, color=MUTED)
    return slide


def slide_section(prs: Presentation, number: str, title: str, subtitle: str, accent: RGBColor = CORAL):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, PAPER)
    # Big number
    add_text(slide, number, left=Inches(1.0), top=Inches(1.6), width=Inches(4), height=Inches(2), size=200, bold=True, color=accent, font_name=FONT_DISPLAY)
    add_text(slide, title, left=Inches(5.5), top=Inches(2.3), width=Inches(7.5), height=Inches(1.2), size=44, bold=True, color=INK, font_name=FONT_DISPLAY)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.45), Inches(1.0), Inches(0.04))
    line.line.fill.background(); line.fill.solid(); line.fill.fore_color.rgb = accent
    add_text(slide, subtitle, left=Inches(5.5), top=Inches(3.65), width=Inches(7.5), height=Inches(2.5), size=18, color=INK_SOFT)
    add_footer(slide, len(prs.slides))
    return slide


def slide_agenda(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, PAPER)
    add_brand_corner(slide)
    add_text(slide, "What you'll see in the next 10 minutes", left=Inches(0.9), top=Inches(0.7), width=Inches(11.5), height=Inches(0.7), size=32, bold=True, color=INK, font_name=FONT_DISPLAY)
    add_text(slide, "A guided tour from problem framing to the live student journey", left=Inches(0.9), top=Inches(1.4), width=Inches(11.5), height=Inches(0.5), size=16, color=MUTED)

    items = [
        ("01", "The Gap We See", "Why students struggle to articulate their skills, and why advisors need leverage", AMBER),
        ("02", "Student Journey", "Upload  →  Skills with Evidence  →  Roles  →  Resume Studio  →  Statement", CORAL),
        ("03", "HK-Local Job Market", "Curriculum-mapped roles for BASc(SDS) and BSc(IM), refreshed monthly", SAGE),
        ("04", "Trust & Governance", "Consent, refusal UX, true deletion, audit  —  PDPO-aligned by design", SKY),
        ("05", "Pilot Plan with CEDARS", "What we propose, with whom, in what scope, on what timeline", AMBER),
    ]

    top = Inches(2.2)
    for i, (num, title, sub, color) in enumerate(items):
        row_top = top + Inches(i * 0.85)
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), row_top, Inches(0.9), Inches(0.65))
        badge.adjustments[0] = 0.4
        badge.line.fill.background(); badge.fill.solid(); badge.fill.fore_color.rgb = color
        tf = badge.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = num; r.font.bold = True; r.font.size = Pt(18); r.font.color.rgb = INK
        add_text(slide, title, left=Inches(2.0), top=row_top + Inches(0.05), width=Inches(4.0), height=Inches(0.5), size=18, bold=True, color=INK)
        add_text(slide, sub, left=Inches(6.1), top=row_top + Inches(0.1), width=Inches(7.0), height=Inches(0.5), size=13, color=INK_SOFT)
    add_footer(slide, len(prs.slides))
    return slide


def slide_split(
    prs: Presentation,
    *,
    eyebrow: str | None,
    title: str,
    bullets: list[str],
    image: Path,
    step: tuple[int, int] | None = None,
    image_first: bool = False,
    note: str | None = None,
):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, PAPER)
    add_brand_corner(slide)

    # Layout regions
    left_x = Inches(0.9) if not image_first else Inches(7.4)
    text_w = Inches(5.6)
    img_x = Inches(7.0) if not image_first else Inches(0.6)
    img_w = Inches(5.7)
    img_h = Inches(5.0)
    img_y = Inches(1.5)

    if eyebrow:
        add_text(slide, eyebrow.upper(), left=left_x, top=Inches(0.65), width=text_w, height=Inches(0.4), size=11, bold=True, color=CORAL)
    add_text(slide, title, left=left_x, top=Inches(1.05), width=text_w, height=Inches(1.3), size=30, bold=True, color=INK, font_name=FONT_DISPLAY, line_spacing=1.05)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_x, Inches(2.55), Inches(0.7), Inches(0.04))
    line.line.fill.background(); line.fill.solid(); line.fill.fore_color.rgb = CORAL

    add_bullets(slide, bullets, left=left_x, top=Inches(2.85), width=text_w, height=Inches(3.4), size=14)

    if note:
        nb = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_x, Inches(6.15), text_w, Inches(0.6))
        nb.adjustments[0] = 0.3
        nb.line.fill.background(); nb.fill.solid(); nb.fill.fore_color.rgb = PAPER_DEEP
        tf = nb.text_frame; tf.margin_left = Emu(120000); tf.margin_right = Emu(120000); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = "▸ "; r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = CORAL
        r2 = p.add_run(); r2.text = note; r2.font.size = Pt(12); r2.font.color.rgb = INK_SOFT

    add_rounded_image(slide, image, left=img_x, top=img_y, width=img_w, height=img_h)

    if step:
        add_step_bar(slide, step[0], step[1], top=Inches(6.85))
    add_footer(slide, len(prs.slides))
    return slide


def slide_hero(prs: Presentation, *, eyebrow: str, title: str, sub: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, PAPER)
    add_brand_corner(slide)
    add_text(slide, eyebrow.upper(), left=Inches(1.0), top=Inches(1.9), width=Inches(11), height=Inches(0.4), size=13, bold=True, color=CORAL, align=PP_ALIGN.CENTER)
    add_text(slide, title, left=Inches(1.0), top=Inches(2.4), width=Inches(11.3), height=Inches(2.6), size=44, bold=True, color=INK, align=PP_ALIGN.CENTER, font_name=FONT_DISPLAY, line_spacing=1.15)
    add_text(slide, sub, left=Inches(1.5), top=Inches(5.1), width=Inches(10.3), height=Inches(1.0), size=18, color=INK_SOFT, align=PP_ALIGN.CENTER)
    add_footer(slide, len(prs.slides))
    return slide


def slide_data_strip(prs: Presentation, *, title: str, sub: str, stats: list[tuple[str, str, str]]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, PAPER)
    add_brand_corner(slide)
    add_text(slide, title, left=Inches(0.9), top=Inches(0.8), width=Inches(11.5), height=Inches(0.8), size=32, bold=True, color=INK, font_name=FONT_DISPLAY)
    add_text(slide, sub, left=Inches(0.9), top=Inches(1.6), width=Inches(11.5), height=Inches(0.5), size=15, color=MUTED)

    n = len(stats)
    card_w = Inches(3.85)
    gap = Inches(0.3)
    total = card_w * n + gap * (n - 1)
    start_left = (SLIDE_W - total) / 2
    top = Inches(2.7)
    for i, (big, label, sub2) in enumerate(stats):
        left = start_left + (card_w + gap) * i
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, Inches(3.4))
        card.adjustments[0] = 0.04
        card.line.fill.background(); card.fill.solid(); card.fill.fore_color.rgb = WHITE
        # Big number
        add_text(slide, big, left=left + Inches(0.3), top=top + Inches(0.4), width=card_w - Inches(0.6), height=Inches(1.3), size=64, bold=True, color=CORAL, align=PP_ALIGN.CENTER, font_name=FONT_DISPLAY)
        add_text(slide, label, left=left + Inches(0.3), top=top + Inches(1.8), width=card_w - Inches(0.6), height=Inches(0.5), size=16, bold=True, color=INK, align=PP_ALIGN.CENTER)
        add_text(slide, sub2, left=left + Inches(0.3), top=top + Inches(2.35), width=card_w - Inches(0.6), height=Inches(0.95), size=12, color=INK_SOFT, align=PP_ALIGN.CENTER)
    add_text(slide, "Source: HKU CEDARS interviews & internal pilot interviews, Mar 2026 (illustrative).", left=Inches(0.9), top=Inches(6.5), width=Inches(11.5), height=Inches(0.3), size=10, color=MUTED, align=PP_ALIGN.CENTER)
    add_footer(slide, len(prs.slides))
    return slide


def slide_flow(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, PAPER)
    add_brand_corner(slide)
    add_text(slide, "How SkillSight works", left=Inches(0.9), top=Inches(0.7), width=Inches(11.5), height=Inches(0.7), size=32, bold=True, color=INK, font_name=FONT_DISPLAY)
    add_text(slide, "From any artefact a student already has, to a verifiable career-readiness statement.", left=Inches(0.9), top=Inches(1.4), width=Inches(11.5), height=Inches(0.5), size=15, color=MUTED)

    inputs = ["Resume\n(PDF / DOCX)", "Project\n(text / code)", "Coursework\n(transcripts)", "Video / Audio\n(presentations)"]
    pipeline = ["Parse\n& Chunk", "Embed\n(vector index)", "Skill\nMatch", "Role\nAlign", "Action\nPlan"]
    outputs = ["Skill Profile\n+ Why / Evidence", "Role Readiness\n+ Gaps", "Recommended\nNext Actions", "Verifiable\nStatement"]

    def draw_node(text, left, top, width, height, fill, txt_color=INK):
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        s.adjustments[0] = 0.25
        s.line.fill.background(); s.fill.solid(); s.fill.fore_color.rgb = fill
        tf = s.text_frame; tf.margin_left = Emu(60000); tf.margin_right = Emu(60000); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        for i, line in enumerate(text.split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = line
            r.font.name = FONT_BODY; r.font.size = Pt(11); r.font.color.rgb = txt_color
            if i == 0:
                r.font.bold = True
        return s

    # Inputs column
    add_text(slide, "INPUTS", left=Inches(0.6), top=Inches(2.2), width=Inches(2.4), height=Inches(0.3), size=10, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
    for i, t in enumerate(inputs):
        draw_node(t, Inches(0.6), Inches(2.55) + Inches(i * 0.95), Inches(2.4), Inches(0.85), AMBER)

    # Pipeline center
    add_text(slide, "ENGINE", left=Inches(3.4), top=Inches(2.2), width=Inches(6.5), height=Inches(0.3), size=10, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
    seg_w = Inches(1.18)
    seg_gap = Inches(0.1)
    for i, t in enumerate(pipeline):
        left = Inches(3.4) + (seg_w + seg_gap) * i
        draw_node(t, left, Inches(4.05), seg_w, Inches(0.95), CORAL, txt_color=WHITE)
        if i < len(pipeline) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, left + seg_w - Inches(0.05), Inches(4.4), Inches(0.18), Inches(0.25))
            arrow.line.fill.background(); arrow.fill.solid(); arrow.fill.fore_color.rgb = INK_SOFT

    # Outputs column
    add_text(slide, "OUTPUTS", left=Inches(10.3), top=Inches(2.2), width=Inches(2.6), height=Inches(0.3), size=10, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
    for i, t in enumerate(outputs):
        draw_node(t, Inches(10.3), Inches(2.55) + Inches(i * 0.95), Inches(2.6), Inches(0.85), SAGE)

    # Caption
    add_text(slide, "Every output carries a back-pointer to the originating evidence chunk(s). No black-box claims.",
             left=Inches(0.9), top=Inches(6.2), width=Inches(11.5), height=Inches(0.4), size=12, color=INK_SOFT, align=PP_ALIGN.CENTER)
    add_footer(slide, len(prs.slides))
    return slide


def slide_image_grid(prs: Presentation, *, eyebrow: str, title: str, sub: str, items: list[tuple[Path, str]]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, PAPER)
    add_brand_corner(slide)
    add_text(slide, eyebrow.upper(), left=Inches(0.9), top=Inches(0.65), width=Inches(11.5), height=Inches(0.35), size=11, bold=True, color=CORAL)
    add_text(slide, title, left=Inches(0.9), top=Inches(1.0), width=Inches(11.5), height=Inches(0.7), size=30, bold=True, color=INK, font_name=FONT_DISPLAY)
    add_text(slide, sub, left=Inches(0.9), top=Inches(1.7), width=Inches(11.5), height=Inches(0.5), size=14, color=INK_SOFT)

    n = len(items)
    cols = min(n, 3)
    rows = (n + cols - 1) // cols
    grid_left = Inches(0.9)
    grid_top = Inches(2.4)
    grid_w = Inches(11.5)
    grid_h = Inches(4.4)
    cell_w = (grid_w - Inches((cols - 1) * 0.25)) / cols
    cell_h = (grid_h - Inches((rows - 1) * 0.25)) / rows
    cap_h = Inches(0.4)
    img_h = cell_h - cap_h - Inches(0.05)
    for i, (img_path, cap) in enumerate(items):
        r, c = divmod(i, cols)
        cell_left = grid_left + (cell_w + Inches(0.25)) * c
        cell_top = grid_top + (cell_h + Inches(0.25)) * r
        add_rounded_image(slide, img_path, left=cell_left, top=cell_top, width=cell_w, height=img_h)
        add_text(slide, cap, left=cell_left, top=cell_top + img_h + Inches(0.05), width=cell_w, height=cap_h, size=12, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_footer(slide, len(prs.slides))
    return slide


def slide_governance(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, PAPER)
    add_brand_corner(slide)
    add_text(slide, "TRUST & GOVERNANCE", left=Inches(0.9), top=Inches(0.65), width=Inches(11.5), height=Inches(0.35), size=11, bold=True, color=CORAL)
    add_text(slide, "Privacy-by-design  ·  PDPO-aligned  ·  Auditable end-to-end", left=Inches(0.9), top=Inches(1.0), width=Inches(11.5), height=Inches(0.7), size=28, bold=True, color=INK, font_name=FONT_DISPLAY)

    pillars = [
        ("Consent up-front", "Every upload requires an explicit purpose (skill assessment / role alignment / portfolio) and scope (full / excerpt / summary). No assumed consent.", AMBER),
        ("Refusal over hallucination", "When evidence is missing, SkillSight refuses to label a skill — it tells the student exactly what to upload instead.", CORAL),
        ("True deletion", "Withdraw consent and we physically delete the file, all extracted chunks, embeddings, assessments and role mappings. Only minimal audit metadata is retained.", SAGE),
        ("Auditable trail", "Every BFF call is recorded with request id, action, status, and subject id. Every change to a skill or role readiness lives in the change log.", SKY),
    ]
    top = Inches(2.0)
    card_w = Inches(5.95)
    card_h = Inches(2.25)
    for i, (title, body, accent) in enumerate(pillars):
        r, c = divmod(i, 2)
        cl = Inches(0.9) + (card_w + Inches(0.3)) * c
        ct = top + (card_h + Inches(0.25)) * r
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cl, ct, card_w, card_h)
        card.adjustments[0] = 0.04
        card.line.fill.background(); card.fill.solid(); card.fill.fore_color.rgb = WHITE
        tab = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cl, ct, Inches(0.18), card_h)
        tab.adjustments[0] = 0.5; tab.line.fill.background(); tab.fill.solid(); tab.fill.fore_color.rgb = accent
        add_text(slide, title, left=cl + Inches(0.45), top=ct + Inches(0.25), width=card_w - Inches(0.6), height=Inches(0.5), size=18, bold=True, color=INK)
        add_text(slide, body, left=cl + Inches(0.45), top=ct + Inches(0.85), width=card_w - Inches(0.6), height=Inches(1.3), size=12, color=INK_SOFT, line_spacing=1.25)
    add_footer(slide, len(prs.slides))
    return slide


def slide_team(prs: Presentation, image: Path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, PAPER)
    add_brand_corner(slide)
    add_text(slide, "WHO BUILT THIS", left=Inches(0.9), top=Inches(0.7), width=Inches(11.5), height=Inches(0.4), size=11, bold=True, color=CORAL)
    add_text(slide, "GEM Lab  ·  Faculty of Education, HKU", left=Inches(0.9), top=Inches(1.05), width=Inches(11.5), height=Inches(0.7), size=30, bold=True, color=INK, font_name=FONT_DISPLAY)
    add_text(slide, "SkillSight is developed by the GEM Lab in the Mathematics, Science, and Technology academic group at HKU's Faculty of Education, led by Prof. Jionghao Lin and his PhD student Ms. Xianghui Meng.",
             left=Inches(0.9), top=Inches(1.85), width=Inches(11.5), height=Inches(1.2), size=15, color=INK_SOFT, line_spacing=1.3)

    add_rounded_image(slide, image, left=Inches(1.5), top=Inches(3.3), width=Inches(10.3), height=Inches(2.8))
    add_text(slide, "Project lead: Prof. Jionghao Lin  ·  https://web.edu.hku.hk/faculty-academics/jionghao",
             left=Inches(0.9), top=Inches(6.35), width=Inches(11.5), height=Inches(0.4), size=13, color=INK, align=PP_ALIGN.CENTER)
    add_footer(slide, len(prs.slides))
    return slide


def slide_pilot(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, PAPER)
    add_brand_corner(slide)
    add_text(slide, "WHAT WE PROPOSE TO CEDARS", left=Inches(0.9), top=Inches(0.7), width=Inches(11.5), height=Inches(0.4), size=11, bold=True, color=CORAL)
    add_text(slide, "A focused pilot on two HKU programmes we already cover", left=Inches(0.9), top=Inches(1.05), width=Inches(11.5), height=Inches(0.8), size=28, bold=True, color=INK, font_name=FONT_DISPLAY)

    add_text(slide, "At this stage, SkillSight is configured for the two HKU programmes whose 2025-26 regulations and syllabuses we have ingested:",
             left=Inches(0.9), top=Inches(1.95), width=Inches(11.5), height=Inches(1.0), size=15, color=INK_SOFT, line_spacing=1.3)

    progs = [
        ("BSc(IM)", "Bachelor of Science in Information Management", "RS 2025-26 (Final)", AMBER),
        ("BASc(SDS)", "Bachelor of Arts and Sciences in Social Data Science", "RS 2025-26 (2025 intake, Final)", SAGE),
    ]
    top = Inches(3.05)
    card_w = Inches(5.95)
    card_h = Inches(2.0)
    for i, (code, full_name, src, accent) in enumerate(progs):
        cl = Inches(0.9) + (card_w + Inches(0.3)) * i
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cl, top, card_w, card_h)
        card.adjustments[0] = 0.05
        card.line.fill.background(); card.fill.solid(); card.fill.fore_color.rgb = WHITE
        tab = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cl, top, Inches(0.18), card_h)
        tab.adjustments[0] = 0.5; tab.line.fill.background(); tab.fill.solid(); tab.fill.fore_color.rgb = accent
        add_text(slide, code, left=cl + Inches(0.45), top=top + Inches(0.2), width=card_w - Inches(0.6), height=Inches(0.6), size=24, bold=True, color=INK, font_name=FONT_DISPLAY)
        add_text(slide, full_name, left=cl + Inches(0.45), top=top + Inches(0.85), width=card_w - Inches(0.6), height=Inches(0.55), size=14, color=INK_SOFT)
        add_text(slide, "Source ingested: " + src, left=cl + Inches(0.45), top=top + Inches(1.45), width=card_w - Inches(0.6), height=Inches(0.4), size=11, color=MUTED)

    add_text(slide, "What's already loaded for these two programmes",
             left=Inches(0.9), top=Inches(5.25), width=Inches(11.5), height=Inches(0.4), size=14, bold=True, color=INK)
    add_bullets(slide, [
        "Curriculum-mapped skill taxonomy (HKU.SKILL.* identifiers)",
        "Potential job catalogue derived from each programme's career outlook + LinkedIn HK snapshots",
        "Ready-to-demo student session, refusal flows, audit log, deletion verification",
    ], left=Inches(0.9), top=Inches(5.65), width=Inches(11.5), height=Inches(1.4), size=13)

    add_footer(slide, len(prs.slides))
    return slide


def slide_closing(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, PAPER)
    blob = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), Inches(-2), Inches(9), Inches(9))
    blob.line.fill.background(); blob.fill.solid(); blob.fill.fore_color.rgb = SKY
    add_text(slide, "Thank you.", left=Inches(0.9), top=Inches(2.4), width=Inches(11.5), height=Inches(1.6), size=88, bold=True, color=INK, font_name=FONT_DISPLAY)
    add_text(slide, "Questions, edge-cases, what would you push us on?", left=Inches(0.9), top=Inches(4.0), width=Inches(11.5), height=Inches(0.8), size=22, color=INK_SOFT)
    add_text(slide, "Live demo  ·  https://skillsight-230.pages.dev/",
             left=Inches(0.9), top=Inches(5.2), width=Inches(11.5), height=Inches(0.5), size=16, color=INK)
    add_text(slide, "GEM Lab, Faculty of Education, The University of Hong Kong",
             left=Inches(0.9), top=Inches(5.65), width=Inches(11.5), height=Inches(0.5), size=14, color=INK_SOFT)
    add_footer(slide, len(prs.slides))
    return slide


# ---------- Build ----------

def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1 — Cover
    slide_cover(prs)

    # 2 — Agenda
    slide_agenda(prs)

    # 3 — Section 1
    slide_section(prs, "01", "Why SkillSight", "Career advisors are stretched thin. Students struggle to articulate what they can do. SkillSight gives both sides a shared, evidence-anchored language.")

    # 4 — The Gap (data strip)
    slide_data_strip(
        prs,
        title="The gap we keep hearing about",
        sub="From conversations with HKU students and CEDARS staff, three patterns repeat:",
        stats=[
            ("3 of 4", "fresh-grad CVs lack evidence", "Most resumes claim skills with no traceable artefact behind them."),
            ("1 : 300+", "advisor-to-student ratio", "Career advisors cannot meaningfully review every portfolio one by one."),
            ("Months", "of lag between JDs & syllabi", "HK job descriptions move faster than course catalogues are revised."),
        ],
    )

    # 5 — Hero one-liner
    slide_hero(
        prs,
        eyebrow="Our answer in one sentence",
        title="Every skill claim, traceable to original evidence.\nEvery gap, paired with the next action.",
        sub="Evidence-First  ·  Explainable  ·  Privacy-by-Design  ·  HK-Local",
    )

    # 6 — How it works (flow)
    slide_flow(prs)

    # 7 — Section 2
    slide_section(prs, "02", "The Student Journey", "What a HKU student does inside SkillSight, end-to-end. Every screenshot is from the live demo at skillsight-230.pages.dev.", accent=AMBER)

    # 8 — Step 1 Login
    slide_split(
        prs,
        eyebrow="Step 1  ·  Sign in",
        title="Frictionless entry, demo-ready",
        bullets=[
            "Three sign-in paths: HKU Portal SSO (production), email link, and a no-friction \u201cTry our demo\u201d for evaluators.",
            "Language toggle in the corner: English / Traditional / Simplified Chinese — every page is fully localised.",
            "The demo session loads with seeded data so CEDARS staff can explore in 30 seconds.",
        ],
        image=SHOTS / "06_login.png",
        step=(1, 8),
    )

    # 9 — Step 2 Dashboard
    slide_split(
        prs,
        eyebrow="Step 2  ·  Dashboard",
        title="One screen tells the student where they stand",
        bullets=[
            "Live counts of authorised documents, verified skills, in-progress assessments and matched roles.",
            "Recommended next steps personalised to the student's current evidence (e.g. \u201cAdd a data project\u201d, \u201cReview my resume\u201d).",
            "Best-matching roles with readiness percentage and surfaced \u201cmust-have\u201d gaps. HKU CEDARS contact card stays visible at the bottom.",
        ],
        image=SHOTS / "07_dashboard.png",
        step=(2, 8),
        image_first=True,
    )

    # 10 — Step 3 Upload Evidence
    slide_split(
        prs,
        eyebrow="Step 3  ·  Upload evidence",
        title="Documents in. Consent collected up-front.",
        bullets=[
            "Multimodal intake: PDFs, DOCX, slides, code, screenshots, video, GitHub repositories.",
            "Per Protocol 9, every upload requires a Purpose (skill assessment / role alignment / portfolio) and a Scope (full / excerpt / summary).",
            "Helpful upload categories — Course Outlines, Assignments, CV / Resume — guide students to high-signal artefacts.",
        ],
        image=SHOTS / "22_evidence.png",
        step=(3, 8),
        note="Refusal-first: if consent is missing or revoked, SkillSight refuses processing and explains the next step.",
    )

    # 11 — Step 4 My Skills
    slide_split(
        prs,
        eyebrow="Step 4  ·  Skills profile",
        title="Click any skill. See the WHY and the EVIDENCE.",
        bullets=[
            "Skill radar gives a quick visual; tabs split skills by status (Verified / Mentioned / Insufficient Evidence / Unassessed).",
            "Expanding a skill reveals the assessment rationale and the original snippet — with a \u201cView Source Location\u201d link back to the chunk.",
            "Insufficient-evidence skills get a refusal banner with a clear next-step prompt; nothing is fabricated.",
        ],
        image=SHOTS / "09b_my_skills_expanded.png",
        step=(4, 8),
        image_first=True,
        note="This is the slide CEDARS will care about most: traceable, refusal-friendly skill claims.",
    )

    # 12 — Step 5 Assessments
    slide_split(
        prs,
        eyebrow="Step 5  ·  Interactive assessments",
        title="Six native assessment types, scored and anchored back to the profile",
        bullets=[
            "Communication, Programming, Writing, Data Analysis, Problem Solving, Presentation — all run inside SkillSight.",
            "Interview Prep Mode generates questions from the student's chosen target role and surfaced gaps.",
            "Every completed assessment becomes a new piece of evidence on the skill profile, with a recorded score.",
        ],
        image=SHOTS / "23_assess.png",
        step=(5, 8),
    )

    # 13 — Step 6 Resume Studio
    slide_split(
        prs,
        eyebrow="Step 6  ·  Resume Enhancement",
        title="A resume coach that suggests — but never overwrites without consent",
        bullets=[
            "Five-step wizard: choose resume → rubric scoring → AI suggestions (accept / reject / edit) → re-score → polished export.",
            "Suggestions ship with a disclaimer: \u201csuggestive, not decisive\u201d. The student stays in the driver's seat.",
            "Targets a specific role; on export students get DOCX or PDF rendered through a vetted template recommended for that role.",
        ],
        image=SHOTS / "11_resume.png",
        step=(6, 8),
        image_first=True,
    )

    # 14 — Step 7 Job Matching
    slide_split(
        prs,
        eyebrow="Step 7  ·  Job matching",
        title="Roles ranked by readiness, gaps named explicitly",
        bullets=[
            "Per-role readiness percentage, paired with the specific skill gaps still standing between the student and \u201cReady\u201d.",
            "\u201cAlmost Ready\u201d vs \u201cReady\u201d status flags help advisors triage who needs which conversation.",
            "Compare-up to 3 roles side-by-side to weigh trade-offs.",
        ],
        image=SHOTS / "12b_jobs_top.png",
        step=(7, 8),
    )

    # 15 — Step 8 Verifiable statement
    slide_split(
        prs,
        eyebrow="Step 8  ·  Verifiable skill statement",
        title="A printable statement students can attach to applications",
        bullets=[
            "Skills, levels and supporting evidence pointers consolidated into a single document.",
            "Print or export as PDF; statements carry generation timestamp and a verify URL for third-party checking.",
            "Refusal-by-default: only verified or demonstrated skills are listed. Insufficient-evidence claims are filtered out.",
        ],
        image=SHOTS / "07_dashboard.png",
        step=(8, 8),
        image_first=True,
        note="Demo URL: /export — the public statement page (PDF render).",
    )

    # 16 — Section 3
    slide_section(prs, "03", "Built for Hong Kong", "We don't recycle generic overseas data. We use HKU-curated taxonomies and live HK job snapshots.", accent=SAGE)

    # 17 — Live HK jobs
    slide_split(
        prs,
        eyebrow="HK-local job pipeline",
        title="Live Jobs — Hong Kong",
        bullets=[
            "Ingests public HK job snapshots (JobsDB, CTgoodjobs and others) into the same skill taxonomy.",
            "Salary bands shown in HKD; role family and source surfaced per posting for transparency.",
            "Matching score is computed against the student's verified profile, not just keyword overlap.",
        ],
        image=SHOTS / "13b_jobs_live_top.png",
    )

    # 18 — Market insights
    slide_split(
        prs,
        eyebrow="HK Market Insights",
        title="Demand signals, not anecdotes",
        bullets=[
            "Skill Demand Trends are computed from real role requirements in the loaded HK postings.",
            "Salary Reference (HKD) gives students realistic anchors; advisors get a shared reference for conversations.",
            "Refresh cadence: monthly snapshots, with full source provenance kept for audit.",
        ],
        image=SHOTS / "16_market.png",
        image_first=True,
    )

    # 19 — Peer Benchmark
    slide_split(
        prs,
        eyebrow="Peer Benchmarking",
        title="\u201cAm I where my classmates are?\u201d",
        bullets=[
            "Each verified skill is positioned as a percentile against an anonymised peer cohort.",
            "No personally-identifiable data is shown — students see only their own placement.",
            "Helps advisors guide priorities: invest where the student is below cohort, defend strengths where they're already top-tier.",
        ],
        image=SHOTS / "17_peer.png",
    )

    # 20 — Section 4
    slide_section(prs, "04", "Trust & Governance", "The reasons CEDARS can endorse SkillSight to the wider University: it is consent-bound, refusal-first and physically deletable.", accent=SKY)

    # 21 — Privacy console
    slide_split(
        prs,
        eyebrow="Privacy console for students",
        title="The full ledger of what we hold, in one screen",
        bullets=[
            "Active vs revoked consents, plus a list of every authorised document with purpose, scope and upload time.",
            "\u201cRevoke & Delete\u201d on any record triggers cascade deletion across DB, vector store and filesystem.",
            "Inline copy explains exactly what gets deleted — original file, chunks, assessments, embeddings, search indexes.",
        ],
        image=SHOTS / "19_privacy.png",
        image_first=True,
    )

    # 22 — Governance pillars
    slide_governance(prs)

    # 23 — Change log
    slide_split(
        prs,
        eyebrow="Audit & change log",
        title="Every skill change, every readiness diff, time-stamped and explained",
        bullets=[
            "Skill status transitions (e.g. not_enough_information → demonstrated) are recorded with timestamps and reason codes.",
            "Role readiness deltas are kept so advisors can see how a student progressed across a semester.",
            "All entries reference canonical HKU.SKILL.* and HKU.ROLE.* identifiers, so the log is interoperable with the registry.",
        ],
        image=SHOTS / "21_change_log.png",
    )

    # 24 — Section 5
    slide_section(prs, "05", "What we propose to CEDARS", "A focused pilot — small enough to run, real enough to learn from.", accent=AMBER)

    # 25 — Pilot proposal
    slide_pilot(prs)

    # 26 — Admin landing (System dashboard)
    slide_split(
        prs,
        eyebrow="For HKU staff",
        title="An admin surface to maintain the taxonomy, jobs and audit trail",
        bullets=[
            "Onboarding (faculties, programmes, courses, terms), Skills registry, Roles library — all editable by HKU staff.",
            "Audit Log: searchable trail of every BFF request, every consent action, every deletion.",
            "Metrics dashboard: usage and reliability; refusal rate as a first-class quality signal.",
        ],
        image=SHOTS / "30_admin.png",
        image_first=True,
        note="Staff and Programme portals are wired into the same canonical registry; aggregated views show no individual student data.",
    )

    # 27 — Team
    slide_team(prs, ASSETS / "team.png")

    # 28 — Closing
    slide_closing(prs)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(f"\nSaved deck: {OUT_PATH}\nSlides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
