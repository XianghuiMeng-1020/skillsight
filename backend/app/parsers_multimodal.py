"""
Multimodal Parsers for SkillSight
- Image OCR (JPG, PNG, WEBP, etc.)
- Video/Audio Transcription
- Extended document formats
"""
import hashlib
import io
import os
import re
import tempfile
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple
import base64


def _make_snippet(text: str, max_len: int = 220) -> str:
    """Create a snippet from text."""
    t = (text or "").strip().replace("\n", " ")
    t = re.sub(r"\s+", " ", t)
    if len(t) <= max_len:
        return t
    return t[:max_len] + "..."


def _compute_hash(text: str) -> str:
    """Compute SHA-256 hash of text."""
    return hashlib.sha256((text or "").encode("utf-8")).hexdigest()


# ====================
# Image Parser (OCR)
# ====================
def parse_image_to_chunks(
    file_path: str = None,
    file_bytes: bytes = None,
    min_chunk_len: int = 30,
) -> List[Dict[str, Any]]:
    """
    Parse image using OCR to extract text.
    Supports: JPG, JPEG, PNG, WEBP, BMP, TIFF (PIL-compatible raster).
    SVG/ICO/HEIC may fail to open and return a placeholder.
    """
    _img_placeholder = [{
        "idx": 0,
        "char_start": 0,
        "char_end": 0,
        "chunk_text": "[Image uploaded - OCR not available or format not supported (e.g. SVG/HEIC). Install pytesseract or easyocr for supported formats: JPG, PNG, BMP, TIFF.]",
        "snippet": "[Image - OCR unavailable]",
        "quote_hash": _compute_hash("image_no_ocr"),
        "section_path": None,
        "page_start": None,
        "page_end": None,
        "media_type": "image",
    }]
    # Try to open image first; PIL does not support SVG, and HEIC needs pillow-heif
    try:
        from PIL import Image
        if file_path:
            img = Image.open(file_path)
        else:
            img = Image.open(io.BytesIO(file_bytes))
        img.load()
    except Exception:
        return _img_placeholder
    # Try pytesseract first
    try:
        import pytesseract
        
        # Perform OCR
        text = pytesseract.image_to_string(img)
        
        if text.strip():
            return _text_to_chunks(text, min_chunk_len, source_type="image_ocr")
        
    except Exception:
        pass
    
    # Try easyocr as fallback
    try:
        import easyocr
        
        reader = easyocr.Reader(['en', 'ch_sim'])  # English + Chinese
        
        if file_path:
            result = reader.readtext(file_path)
        else:
            # Save to temp file for easyocr
            with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                if file_bytes:
                    tmp.write(file_bytes)
                else:
                    tmp.write(Path(file_path).read_bytes())
                tmp_path = tmp.name
            try:
                result = reader.readtext(tmp_path)
            finally:
                os.unlink(tmp_path)
        
        # Combine detected text
        text = " ".join([r[1] for r in result])
        
        if text.strip():
            return _text_to_chunks(text, min_chunk_len, source_type="image_ocr")
        return _img_placeholder
        
    except Exception:
        pass
    
    return _img_placeholder


def _detect_image_type(data: bytes) -> str:
    """Detect image type from file header bytes (replaces removed imghdr)."""
    if data[:8] == b'\x89PNG\r\n\x1a\n':
        return "png"
    if data[:2] == b'\xff\xd8':
        return "jpeg"
    if data[:4] == b'GIF8':
        return "gif"
    if data[:4] == b'RIFF' and data[8:12] == b'WEBP':
        return "webp"
    if data[:2] in (b'BM',):
        return "bmp"
    if data[:4] in (b'II\x2a\x00', b'MM\x00\x2a'):
        return "tiff"
    return "jpeg"


def parse_image_for_vision_model(
    file_path: str = None,
    file_bytes: bytes = None,
) -> Optional[Dict[str, Any]]:
    """
    Prepare image for vision-capable LLM (e.g., GPT-4V, Claude 3).
    Returns base64 encoded image and metadata, or None if format unsupported.
    """
    try:
        if file_path:
            with open(file_path, "rb") as f:
                file_bytes = f.read()
        if not file_bytes:
            return None
        img_type = _detect_image_type(file_bytes)
        mime_type = f"image/{img_type}"
        b64_data = base64.b64encode(file_bytes).decode("utf-8")
        return {
            "type": "image",
            "mime_type": mime_type,
            "base64": b64_data,
            "size_bytes": len(file_bytes),
        }
    except Exception:
        return None


# ====================
# Video/Audio Parser (Transcription)
# ====================
def parse_video_to_chunks(
    file_path: str = None,
    file_bytes: bytes = None,
    min_chunk_len: int = 50,
) -> List[Dict[str, Any]]:
    """
    Parse video/audio file to extract transcript.
    Supports: MP4, MP3, WAV, M4A, WEBM, OGG
    
    Uses:
    1. OpenAI Whisper (local or API)
    2. Falls back to placeholder if unavailable
    """
    # Save bytes to temp file if needed
    if file_bytes and not file_path:
        with tempfile.NamedTemporaryFile(suffix=".mp4", delete=False) as tmp:
            tmp.write(file_bytes)
            file_path = tmp.name
            cleanup_temp = True
    else:
        cleanup_temp = False
    
    transcript_text = None
    timestamps = []
    
    # Try local Whisper
    try:
        import whisper
        
        model = whisper.load_model("base")  # Use "small" or "medium" for better accuracy
        result = model.transcribe(file_path)
        
        transcript_text = result["text"]
        
        # Extract timestamps for segments
        for seg in result.get("segments", []):
            timestamps.append({
                "start": seg["start"],
                "end": seg["end"],
                "text": seg["text"],
            })
            
    except ImportError:
        pass
    
    # Try OpenAI Whisper API as fallback
    if not transcript_text:
        try:
            import openai
            
            with open(file_path, "rb") as audio_file:
                result = openai.Audio.transcribe("whisper-1", audio_file)
            transcript_text = result.get("text", "")
            
        except Exception:
            pass
    
    # Cleanup temp file
    if cleanup_temp and os.path.exists(file_path):
        os.unlink(file_path)
    
    if transcript_text and transcript_text.strip():
        chunks = _text_to_chunks(transcript_text, min_chunk_len, source_type="video_transcript")
        
        # Add timestamp info if available
        if timestamps:
            for i, chunk in enumerate(chunks):
                chunk["timestamps"] = timestamps
                chunk["media_type"] = "video"
        
        return chunks
    
    # No transcription available
    return [{
        "idx": 0,
        "char_start": 0,
        "char_end": 0,
        "chunk_text": "[Video/Audio uploaded - Transcription not available. Install openai-whisper for transcription.]",
        "snippet": "[Video/Audio - transcription unavailable]",
        "quote_hash": _compute_hash("video_no_transcript"),
        "section_path": None,
        "page_start": None,
        "page_end": None,
        "media_type": "video",
    }]


def extract_video_metadata(file_path: str) -> Dict[str, Any]:
    """Extract metadata from video file."""
    try:
        import cv2
        
        cap = cv2.VideoCapture(file_path)
        
        metadata = {
            "duration_seconds": cap.get(cv2.CAP_PROP_FRAME_COUNT) / cap.get(cv2.CAP_PROP_FPS),
            "fps": cap.get(cv2.CAP_PROP_FPS),
            "width": int(cap.get(cv2.CAP_PROP_FRAME_WIDTH)),
            "height": int(cap.get(cv2.CAP_PROP_FRAME_HEIGHT)),
            "frame_count": int(cap.get(cv2.CAP_PROP_FRAME_COUNT)),
        }
        
        cap.release()
        return metadata
        
    except Exception:
        return {}


# ====================
# Legacy Word (.doc) Parser
# ====================
def parse_doc_to_chunks(
    file_path: str = None,
    file_bytes: bytes = None,
    min_chunk_len: int = 30,
) -> List[Dict[str, Any]]:
    """
    Parse legacy .doc files (Word 97-2003 format).
    Uses antiword or textract as fallback.
    """
    # Save bytes to temp file if needed
    if file_bytes and not file_path:
        with tempfile.NamedTemporaryFile(suffix=".doc", delete=False) as tmp:
            tmp.write(file_bytes)
            file_path = tmp.name
            cleanup_temp = True
    else:
        cleanup_temp = False
    
    text = None
    
    # Try antiword first (faster, cleaner output)
    try:
        import subprocess
        result = subprocess.run(
            ["antiword", file_path],
            capture_output=True,
            text=True,
            timeout=30
        )
        if result.returncode == 0:
            text = result.stdout
    except Exception:
        pass
    
    # Try textract as fallback
    if not text:
        try:
            import textract
            text = textract.process(file_path).decode("utf-8", errors="ignore")
        except ImportError:
            pass
        except Exception:
            pass
    
    # Try python-docx2txt as another fallback
    if not text:
        try:
            import docx2txt
            text = docx2txt.process(file_path)
        except ImportError:
            pass
        except Exception:
            pass
    
    # Cleanup temp file
    if cleanup_temp and os.path.exists(file_path):
        os.unlink(file_path)
    
    if text and text.strip():
        return _text_to_chunks(text, min_chunk_len, source_type="doc")
    
    # Return placeholder if no parser available
    return [{
        "idx": 0,
        "char_start": 0,
        "char_end": 0,
        "chunk_text": "[Legacy Word .doc file - Install antiword or textract for text extraction]",
        "snippet": "[.doc file - parser unavailable]",
        "quote_hash": _compute_hash("doc_no_parser"),
        "section_path": None,
        "page_start": None,
        "page_end": None,
        "media_type": "document",
    }]


# ====================
# Jupyter Notebook Parser
# ====================
def parse_ipynb_to_chunks(
    file_path: str = None,
    file_bytes: bytes = None,
    min_chunk_len: int = 30,
) -> List[Dict[str, Any]]:
    """
    Parse Jupyter Notebook (.ipynb) files.
    Extracts code cells and markdown cells.
    """
    import json as json_lib
    
    try:
        if file_bytes:
            content = file_bytes.decode("utf-8", errors="ignore")
        else:
            content = Path(file_path).read_text(encoding="utf-8", errors="ignore")
        
        notebook = json_lib.loads(content)
    except Exception as e:
        return [{
            "idx": 0,
            "char_start": 0,
            "char_end": 0,
            "chunk_text": f"[Failed to parse Jupyter notebook: {e}]",
            "snippet": "[Notebook parse error]",
            "quote_hash": _compute_hash("ipynb_error"),
            "section_path": None,
            "page_start": None,
            "page_end": None,
            "media_type": "notebook",
        }]
    
    chunks = []
    idx = 0
    char_cursor = 0
    
    cells = notebook.get("cells", [])
    
    for cell_num, cell in enumerate(cells):
        cell_type = cell.get("cell_type", "")
        source = cell.get("source", [])
        
        # Handle source as list or string
        if isinstance(source, list):
            cell_content = "".join(source)
        else:
            cell_content = str(source)
        
        cell_content = cell_content.strip()
        
        if len(cell_content) < min_chunk_len:
            char_cursor += len(cell_content) + 1
            continue
        
        # Determine section path based on cell type
        if cell_type == "markdown":
            # Try to extract heading from markdown
            lines = cell_content.split("\n")
            heading = None
            for line in lines:
                if line.startswith("#"):
                    heading = line.lstrip("#").strip()
                    break
            section_path = f"Cell {cell_num + 1} (Markdown)" + (f": {heading}" if heading else "")
        elif cell_type == "code":
            # Try to detect function/class definition
            func_match = re.search(r'^(?:def |class |async def )(\w+)', cell_content, re.MULTILINE)
            if func_match:
                section_path = f"Cell {cell_num + 1} (Code): {func_match.group(1)}"
            else:
                section_path = f"Cell {cell_num + 1} (Code)"
        else:
            section_path = f"Cell {cell_num + 1} ({cell_type})"
        
        # Also capture outputs for code cells
        outputs_text = ""
        if cell_type == "code":
            outputs = cell.get("outputs", [])
            for output in outputs:
                if output.get("output_type") == "stream":
                    text = output.get("text", [])
                    if isinstance(text, list):
                        outputs_text += "".join(text)
                    else:
                        outputs_text += str(text)
                elif output.get("output_type") in ("execute_result", "display_data"):
                    data = output.get("data", {})
                    if "text/plain" in data:
                        plain = data["text/plain"]
                        if isinstance(plain, list):
                            outputs_text += "".join(plain)
                        else:
                            outputs_text += str(plain)
        
        # Combine cell content with output (if any)
        full_content = cell_content
        if outputs_text.strip():
            full_content += f"\n\n# Output:\n{outputs_text.strip()}"
        
        chunks.append({
            "idx": idx,
            "char_start": char_cursor,
            "char_end": char_cursor + len(full_content),
            "chunk_text": full_content,
            "snippet": _make_snippet(full_content, 200),
            "quote_hash": _compute_hash(full_content),
            "section_path": section_path,
            "page_start": cell_num + 1,
            "page_end": cell_num + 1,
            "media_type": "notebook",
            "cell_type": cell_type,
        })
        
        char_cursor += len(full_content) + 1
        idx += 1
    
    return chunks if chunks else [{
        "idx": 0,
        "char_start": 0,
        "char_end": 0,
        "chunk_text": "[Empty Jupyter notebook]",
        "snippet": "[Empty notebook]",
        "quote_hash": _compute_hash("ipynb_empty"),
        "section_path": None,
        "page_start": None,
        "page_end": None,
        "media_type": "notebook",
    }]


# ====================
# Spreadsheet Parser
# ====================
def parse_spreadsheet_to_chunks(
    file_path: str = None,
    file_bytes: bytes = None,
    filename: str = None,
    min_chunk_len: int = 30,
) -> List[Dict[str, Any]]:
    """
    Parse Excel spreadsheets (.xlsx, .xls) to extract cell data.
    """
    # Determine format
    if filename:
        ext = os.path.splitext(filename)[1].lower()
    elif file_path:
        ext = os.path.splitext(file_path)[1].lower()
    else:
        ext = ".xlsx"
    
    # Try openpyxl for xlsx
    if ext == ".xlsx":
        try:
            from openpyxl import load_workbook
            
            if file_bytes:
                wb = load_workbook(io.BytesIO(file_bytes), data_only=True)
            else:
                wb = load_workbook(file_path, data_only=True)
            
            chunks = []
            idx = 0
            char_cursor = 0
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                rows_text = []
                
                for row in sheet.iter_rows(values_only=True):
                    # Skip empty rows
                    if all(cell is None for cell in row):
                        continue
                    row_values = [str(cell) if cell is not None else "" for cell in row]
                    rows_text.append(" | ".join(row_values))
                
                if rows_text:
                    sheet_content = f"Sheet: {sheet_name}\n" + "\n".join(rows_text)
                    
                    if len(sheet_content) >= min_chunk_len:
                        chunks.append({
                            "idx": idx,
                            "char_start": char_cursor,
                            "char_end": char_cursor + len(sheet_content),
                            "chunk_text": sheet_content,
                            "snippet": _make_snippet(sheet_content, 200),
                            "quote_hash": _compute_hash(sheet_content),
                            "section_path": f"Sheet: {sheet_name}",
                            "page_start": None,
                            "page_end": None,
                            "media_type": "spreadsheet",
                        })
                        char_cursor += len(sheet_content) + 1
                        idx += 1
            
            return chunks if chunks else [{
                "idx": 0,
                "char_start": 0,
                "char_end": 0,
                "chunk_text": "[Empty spreadsheet]",
                "snippet": "[Empty spreadsheet]",
                "quote_hash": _compute_hash("xlsx_empty"),
                "section_path": None,
                "page_start": None,
                "page_end": None,
                "media_type": "spreadsheet",
            }]
            
        except ImportError:
            pass
        except Exception as e:
            return [{
                "idx": 0,
                "char_start": 0,
                "char_end": 0,
                "chunk_text": f"[Failed to parse spreadsheet: {e}]",
                "snippet": "[Spreadsheet parse error]",
                "quote_hash": _compute_hash("xlsx_error"),
                "section_path": None,
                "page_start": None,
                "page_end": None,
                "media_type": "spreadsheet",
            }]
    
    # Try xlrd for xls
    if ext == ".xls":
        try:
            import xlrd
            
            if file_bytes:
                wb = xlrd.open_workbook(file_contents=file_bytes)
            else:
                wb = xlrd.open_workbook(file_path)
            
            chunks = []
            idx = 0
            char_cursor = 0
            
            for sheet in wb.sheets():
                rows_text = []
                for row_idx in range(sheet.nrows):
                    row = sheet.row_values(row_idx)
                    row_values = [str(cell) if cell else "" for cell in row]
                    if any(row_values):
                        rows_text.append(" | ".join(row_values))
                
                if rows_text:
                    sheet_content = f"Sheet: {sheet.name}\n" + "\n".join(rows_text)
                    
                    if len(sheet_content) >= min_chunk_len:
                        chunks.append({
                            "idx": idx,
                            "char_start": char_cursor,
                            "char_end": char_cursor + len(sheet_content),
                            "chunk_text": sheet_content,
                            "snippet": _make_snippet(sheet_content, 200),
                            "quote_hash": _compute_hash(sheet_content),
                            "section_path": f"Sheet: {sheet.name}",
                            "page_start": None,
                            "page_end": None,
                            "media_type": "spreadsheet",
                        })
                        char_cursor += len(sheet_content) + 1
                        idx += 1
            
            return chunks if chunks else [{
                "idx": 0,
                "char_start": 0,
                "char_end": 0,
                "chunk_text": "[Empty spreadsheet]",
                "snippet": "[Empty spreadsheet]",
                "quote_hash": _compute_hash("xls_empty"),
                "section_path": None,
                "page_start": None,
                "page_end": None,
                "media_type": "spreadsheet",
            }]
            
        except ImportError:
            pass
        except Exception:
            pass
    
    # Fallback
    return [{
        "idx": 0,
        "char_start": 0,
        "char_end": 0,
        "chunk_text": "[Spreadsheet file - Install openpyxl (xlsx) or xlrd (xls) for parsing]",
        "snippet": "[Spreadsheet - parser unavailable]",
        "quote_hash": _compute_hash("spreadsheet_no_parser"),
        "section_path": None,
        "page_start": None,
        "page_end": None,
        "media_type": "spreadsheet",
    }]


# ====================
# Presentation Parser
# ====================
def parse_pptx_to_chunks(
    file_path: str = None,
    file_bytes: bytes = None,
    min_chunk_len: int = 30,
    filename: str = None,
) -> List[Dict[str, Any]]:
    """
    Parse PowerPoint presentation to extract text and notes.
    Only .pptx is supported; .ppt (legacy binary) returns a placeholder.
    """
    ext = ""
    if filename:
        ext = os.path.splitext(filename)[1].lower()
    elif file_path:
        ext = os.path.splitext(file_path)[1].lower()
    if ext == ".ppt":
        return [{
            "idx": 0,
            "char_start": 0,
            "char_end": 0,
            "chunk_text": "[Legacy .ppt format - not supported. Save as .pptx to extract text.]",
            "snippet": "[.ppt - use .pptx]",
            "quote_hash": _compute_hash("ppt_unsupported"),
            "section_path": None,
            "page_start": None,
            "page_end": None,
            "media_type": "presentation",
        }]
    try:
        from pptx import Presentation
        
        if file_bytes:
            prs = Presentation(io.BytesIO(file_bytes))
        else:
            prs = Presentation(file_path)
        
        chunks = []
        idx = 0
        char_cursor = 0
        
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text_parts = []
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text_parts.append(shape.text.strip())
            
            # Extract speaker notes
            notes_text = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
            
            # Combine slide content
            slide_content = "\n".join(slide_text_parts)
            
            if slide_content and len(slide_content) >= min_chunk_len:
                chunks.append({
                    "idx": idx,
                    "char_start": char_cursor,
                    "char_end": char_cursor + len(slide_content),
                    "chunk_text": slide_content,
                    "snippet": _make_snippet(slide_content),
                    "quote_hash": _compute_hash(slide_content),
                    "section_path": f"Slide {slide_num}",
                    "page_start": slide_num,
                    "page_end": slide_num,
                    "media_type": "presentation",
                    "speaker_notes": notes_text if notes_text else None,
                })
                char_cursor += len(slide_content) + 1
                idx += 1
        
        return chunks
        
    except ImportError:
        return [{
            "idx": 0,
            "char_start": 0,
            "char_end": 0,
            "chunk_text": "[PPTX parsing requires python-pptx. Install with: pip install python-pptx. Content stored for reference.]",
            "snippet": "[.pptx - install python-pptx]",
            "quote_hash": _compute_hash("pptx_no_module"),
            "section_path": None,
            "page_start": None,
            "page_end": None,
            "media_type": "presentation",
        }]
    except Exception as e:
        return [{
            "idx": 0,
            "char_start": 0,
            "char_end": 0,
            "chunk_text": f"[Presentation file could not be parsed: {e}. Ensure file is valid .pptx.]",
            "snippet": "[Parse error]",
            "quote_hash": _compute_hash("pptx_error"),
            "section_path": None,
            "page_start": None,
            "page_end": None,
            "media_type": "presentation",
        }]


# ====================
# Helper Functions
# ====================
def _placeholder_chunk(media_type: str, message: str) -> Dict[str, Any]:
    """Return a single placeholder chunk when parsing is not available."""
    return {
        "idx": 0,
        "char_start": 0,
        "char_end": 0,
        "chunk_text": message,
        "snippet": message[:150] + ("..." if len(message) > 150 else ""),
        "quote_hash": _compute_hash(message),
        "section_path": None,
        "page_start": None,
        "page_end": None,
        "media_type": media_type,
    }


def _text_to_chunks(
    text: str,
    min_chunk_len: int = 50,
    source_type: str = "text",
) -> List[Dict[str, Any]]:
    """Convert raw text to chunk format."""
    if not text:
        return []
    
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    chunks = []
    cursor = 0
    idx = 0
    
    for part in text.split("\n\n"):
        part_strip = part.strip()
        if len(part_strip) < min_chunk_len:
            cursor += len(part) + 2
            continue
        
        pos = text.find(part, cursor)
        if pos == -1:
            pos = cursor
        
        char_start = pos
        char_end = pos + len(part)
        
        chunks.append({
            "idx": idx,
            "char_start": char_start,
            "char_end": char_end,
            "chunk_text": part_strip,
            "snippet": _make_snippet(part_strip),
            "quote_hash": _compute_hash(part_strip),
            "section_path": None,
            "page_start": None,
            "page_end": None,
            "source_type": source_type,
        })
        
        cursor = char_end + 2
        idx += 1
    
    return chunks


# ====================
# Unified Multimodal Parser
# ====================
SUPPORTED_EXTENSIONS = {
    # Documents
    ".txt": "text",
    ".doc": "document",  # Legacy Word format
    ".docx": "document",
    ".pdf": "document",
    ".pptx": "presentation",
    ".ppt": "presentation",
    ".rtf": "document",
    ".odt": "document",  # OpenDocument Text
    ".md": "text",  # Markdown
    ".markdown": "text",
    # Spreadsheets
    ".xlsx": "spreadsheet",
    ".xls": "spreadsheet",
    ".csv": "text",
    # Images
    ".jpg": "image",
    ".jpeg": "image",
    ".png": "image",
    ".webp": "image",
    ".bmp": "image",
    ".tiff": "image",
    ".tif": "image",
    ".gif": "image",
    ".svg": "image",
    ".ico": "image",
    ".heic": "image",
    ".heif": "image",
    # Video/Audio
    ".mp4": "video",
    ".webm": "video",
    ".mov": "video",
    ".avi": "video",
    ".mkv": "video",
    ".flv": "video",
    ".wmv": "video",
    ".mp3": "audio",
    ".wav": "audio",
    ".m4a": "audio",
    ".ogg": "audio",
    ".flac": "audio",
    ".aac": "audio",
    # Code - Python
    ".py": "code",
    ".pyw": "code",
    ".pyi": "code",
    ".ipynb": "notebook",  # Jupyter Notebook
    # Code - JavaScript/TypeScript
    ".js": "code",
    ".jsx": "code",
    ".ts": "code",
    ".tsx": "code",
    ".mjs": "code",
    ".cjs": "code",
    ".vue": "code",
    ".svelte": "code",
    # Code - Web
    ".html": "code",
    ".htm": "code",
    ".css": "code",
    ".scss": "code",
    ".sass": "code",
    ".less": "code",
    # Code - Systems
    ".java": "code",
    ".cpp": "code",
    ".cc": "code",
    ".cxx": "code",
    ".c": "code",
    ".h": "code",
    ".hpp": "code",
    ".cs": "code",  # C#
    ".go": "code",
    ".rs": "code",
    ".rb": "code",
    ".php": "code",
    ".swift": "code",
    ".kt": "code",
    ".kts": "code",
    ".scala": "code",
    ".r": "code",
    ".R": "code",
    ".m": "code",  # Objective-C / MATLAB
    ".mm": "code",  # Objective-C++
    # Code - Shell/Scripts
    ".sh": "code",
    ".bash": "code",
    ".zsh": "code",
    ".fish": "code",
    ".ps1": "code",  # PowerShell
    ".bat": "code",
    ".cmd": "code",
    # Code - Data/Config
    ".json": "code",
    ".yaml": "code",
    ".yml": "code",
    ".xml": "code",
    ".toml": "code",
    ".ini": "code",
    ".cfg": "code",
    ".conf": "code",
    ".env": "code",
    # Code - Database
    ".sql": "code",
    # Code - Other
    ".lua": "code",
    ".pl": "code",
    ".pm": "code",
    ".ex": "code",
    ".exs": "code",
    ".erl": "code",
    ".hrl": "code",
    ".clj": "code",
    ".cljs": "code",
    ".hs": "code",
    ".lhs": "code",
    ".elm": "code",
    ".dart": "code",
    ".groovy": "code",
    ".gradle": "code",
    ".tf": "code",  # Terraform
    ".proto": "code",  # Protocol Buffers
    ".graphql": "code",
    ".gql": "code",
    # Misc
    ".log": "text",
    ".diff": "code",
    ".patch": "code",
}


def parse_multimodal_file(
    file_path: str = None,
    file_bytes: bytes = None,
    filename: str = None,
    min_chunk_len: int = 50,
) -> Dict[str, Any]:
    """
    Unified multimodal parser that handles all supported file types.
    
    Returns:
        {
            "chunks": List of chunk dicts,
            "media_type": str (text/image/video/audio/code/presentation),
            "metadata": Dict with file-specific metadata,
            "vision_data": Optional base64 image for vision models,
        }
    """
    try:
        from backend.app.parsers import parse_file_to_chunks, parse_txt_to_chunks
    except ImportError:
        from app.parsers import parse_file_to_chunks, parse_txt_to_chunks
    
    if file_path:
        ext = os.path.splitext(file_path)[1].lower()
        if not filename:
            filename = os.path.basename(file_path)
    elif filename:
        ext = os.path.splitext(filename)[1].lower()
    else:
        raise ValueError("Either file_path or filename must be provided")
    
    media_type = SUPPORTED_EXTENSIONS.get(ext, "unknown")
    result = {
        "chunks": [],
        "media_type": media_type,
        "metadata": {"extension": ext, "filename": filename},
        "vision_data": None,
    }
    
    # Route to appropriate parser
    if ext == ".doc":
        # Legacy Word format - use special parser
        result["chunks"] = parse_doc_to_chunks(
            file_path=file_path,
            file_bytes=file_bytes,
            min_chunk_len=min_chunk_len,
        )
    
    elif ext in (".rtf", ".odt"):
        # RTF/ODT: parsers.py does not support these; try raw text extraction or placeholder
        try:
            raw = file_bytes if file_bytes else (Path(file_path).read_bytes() if file_path else b"")
            text = raw.decode("utf-8", errors="ignore")
            if text.strip():
                result["chunks"] = _text_to_chunks(text, min_chunk_len, source_type="document")
            else:
                result["chunks"] = [_placeholder_chunk("document", "RTF/ODT file - text extraction limited. Content stored for reference.")]
        except Exception:
            result["chunks"] = [_placeholder_chunk("document", "RTF/ODT file - could not extract text. Install optional libraries for full support.")]
    
    elif media_type == "text" or media_type == "document":
        # Use existing document parsers for txt, docx, pdf, md, csv
        try:
            result["chunks"] = parse_file_to_chunks(
                file_path=file_path,
                file_bytes=file_bytes,
                filename=filename,
                min_chunk_len=min_chunk_len,
            )
        except Exception as e:
            result["error"] = str(e)
    
    elif media_type == "notebook":
        # Jupyter Notebook
        result["chunks"] = parse_ipynb_to_chunks(
            file_path=file_path,
            file_bytes=file_bytes,
            min_chunk_len=min_chunk_len,
        )
    
    elif media_type == "spreadsheet":
        # Excel/CSV files
        result["chunks"] = parse_spreadsheet_to_chunks(
            file_path=file_path,
            file_bytes=file_bytes,
            filename=filename,
            min_chunk_len=min_chunk_len,
        )
    
    elif media_type == "presentation":
        result["chunks"] = parse_pptx_to_chunks(
            file_path=file_path,
            file_bytes=file_bytes,
            filename=filename,
            min_chunk_len=min_chunk_len,
        )
    
    elif media_type == "image":
        result["chunks"] = parse_image_to_chunks(
            file_path=file_path,
            file_bytes=file_bytes,
            min_chunk_len=min_chunk_len,
        )
        # Also prepare for vision model
        result["vision_data"] = parse_image_for_vision_model(
            file_path=file_path,
            file_bytes=file_bytes,
        )
    
    elif media_type in ("video", "audio"):
        result["chunks"] = parse_video_to_chunks(
            file_path=file_path,
            file_bytes=file_bytes,
            min_chunk_len=min_chunk_len,
        )
        # Get video metadata
        if file_path and media_type == "video":
            result["metadata"].update(extract_video_metadata(file_path))
    
    elif media_type == "code":
        # Parse code files as text with syntax awareness
        if file_bytes:
            content = file_bytes.decode("utf-8", errors="ignore")
        else:
            content = Path(file_path).read_text(encoding="utf-8", errors="ignore")
        
        # Split by function/class definitions for better chunking
        result["chunks"] = _parse_code_to_chunks(content, ext, min_chunk_len)
    
    else:
        result["error"] = f"Unsupported file type: {ext}"
    
    return result


def _parse_code_to_chunks(
    content: str,
    ext: str,
    min_chunk_len: int = 30,
) -> List[Dict[str, Any]]:
    """Parse code file into logical chunks."""
    chunks = []
    idx = 0
    
    # Simple chunking by double newlines (works for most languages)
    # A more sophisticated approach would use AST parsing
    parts = re.split(r'\n\n+', content)
    char_cursor = 0
    
    for part in parts:
        part_strip = part.strip()
        if len(part_strip) < min_chunk_len:
            char_cursor += len(part) + 2
            continue
        
        # Detect if this is a function/class definition
        section_path = None
        if re.match(r'^(def |class |function |const |let |var |public |private )', part_strip):
            # Extract name
            match = re.match(r'^(?:def |class |function |const |let |var |public |private )(\w+)', part_strip)
            if match:
                section_path = match.group(1)
        
        chunks.append({
            "idx": idx,
            "char_start": char_cursor,
            "char_end": char_cursor + len(part),
            "chunk_text": part_strip,
            "snippet": _make_snippet(part_strip, 150),
            "quote_hash": _compute_hash(part_strip),
            "section_path": section_path,
            "page_start": None,
            "page_end": None,
            "media_type": "code",
            "language": ext.lstrip("."),
        })
        
        char_cursor += len(part) + 2
        idx += 1
    
    return chunks
